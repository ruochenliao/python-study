import base64
import os
import subprocess
from io import BytesIO
# pip install pymupdf
import fitz
from PIL import Image
from pptx import Presentation
from ai_rag_system_v1.utils import settings
from .ocr import ocr_to_text_from_llm


def get_b64_image_from_path(image_path):
  with open(image_path, "rb") as image_file:
    return base64.b64encode(image_file.read()).decode('utf-8')

def process_table(file):
    """
    解析表格:
    1、借助ocr直接识别表格内容
    2、借助多模态大模型识别
    :param file:
    :return: 表格内容及对表格的描述信息
    """
    content = ocr_to_text_from_llm(file)
    llm = settings.deepseek_llm()
    response = llm.complete(f"你的职责是解释表格。"
                            f"你是将线性化表格转换成简单中文文本供大型语言模型（LLMs）使用的专家。"
                            f"请解释以下线性化表格： {content}")
    return content, response.text

def get_b64_image_from_content(image_content):
    """Convert image content to base64 encoded string."""
    img = Image.open(BytesIO(image_content))
    if img.mode != 'RGB':
        img = img.convert('RGB')
    buffered = BytesIO()
    img.save(buffered, format="JPEG")
    return base64.b64encode(buffered.getvalue()).decode("utf-8")

def describe_image(file_path, prompt:str="请尽可能详细的描述你在图片中看到的所有内容"):
    """Generate a description of an image using VLLM API."""
    image_b64 = get_b64_image_from_path(file_path)
    client = settings.vllm()
    messages = [
        {
            "role": "user",
            "content":
                [
                    {
                        "type": "text",
                        "text": prompt
                    },
                    {
                        "type": "image_url",
                        "image_url":{"url": f"data:image/png;base64,{image_b64}"}
                     }
                ]
         }
    ]

    # 调用LLM生成几个相关的问题，多角度获取图片中的内容，自行完成
    # 参考LLmaindex结构化输出，生成多个相关问题
    # https://docs.llamaindex.ai/en/stable/examples/output_parsing/llm_program/

    completion = client.chat.completions.create(
        model=settings.configuration.vllm_model_name,
        messages=messages, temperature=0.5,
        seed=0, top_p=0.70, stream=False
    )
    # 返回多个问题汇总后的答案描述
    return completion.choices[0].message.content

def extract_text_around_item(text_blocks, bbox, page_height, threshold_percentage=0.1):
    """从页面上的给定边界框提取上方和下方的文本。"""
    before_text, after_text = "", ""  # 初始化上方和下方文本为空字符串
    vertical_threshold_distance = page_height * threshold_percentage  # 计算垂直阈值距离
    horizontal_threshold_distance = bbox.width * threshold_percentage  # 计算水平阈值距离

    for block in text_blocks:  # 遍历所有文本块
        block_bbox = fitz.Rect(block[:4])  # 获取当前文本块的边界框
        vertical_distance = min(abs(block_bbox.y1 - bbox.y0), abs(block_bbox.y0 - bbox.y1))  # 计算当前文本块与目标边界框的垂直距离
        horizontal_overlap = max(0, min(block_bbox.x1, bbox.x1) - max(block_bbox.x0, bbox.x0))  # 计算当前文本块与目标边界框的水平重叠

        if vertical_distance <= vertical_threshold_distance and horizontal_overlap >= -horizontal_threshold_distance:
            # 如果垂直距离小于等于阈值且水平重叠大于等于负阈值
            if block_bbox.y1 < bbox.y0 and not before_text:
                # 如果当前文本块在目标边界框上方且上方文本未被设置
                before_text = block[4]  # 更新上方文本
            elif block_bbox.y0 > bbox.y1 and not after_text:
                # 如果当前文本块在目标边界框下方且下方文本未被设置
                after_text = block[4]  # 更新下方文本
                break  # 结束循环

    return before_text, after_text  # 返回提取到的上方和下方文本

def process_text_blocks(text_blocks, char_count_threshold=500):
    """根据字符数阈值对文本块进行分组。"""
    current_group = []  # 当前组的文本块列表
    grouped_blocks = []  # 分组后的文本块列表
    current_char_count = 0  # 当前组的总字符数

    for block in text_blocks:
        if block[-1] == 0:  # 检查块是否为文本类型
            block_text = block[4]  # 获取块的文本内容
            block_char_count = len(block_text)  # 计算块的字符数

            if current_char_count + block_char_count <= char_count_threshold:
                # 如果当前组的字符数加上新块的字符数不超过阈值
                current_group.append(block)  # 将新块添加到当前组
                current_char_count += block_char_count  # 更新当前组的字符计数
            else:
                # 如果超过阈值
                if current_group:
                    # 如果当前组不为空
                    grouped_content = "\n".join([b[4] for b in current_group])  # 合并当前组的内容
                    grouped_blocks.append((current_group[0], grouped_content))  # 将合并后的内容添加到分组后的块列表
                current_group = [block]  # 重置当前组
                current_char_count = block_char_count  # 重置当前组的字符计数

    # 处理最后一个组
    if current_group:
        grouped_content = "\n".join([b[4] for b in current_group])  # 合并最后一个组的内容
        grouped_blocks.append((current_group[0], grouped_content))  # 将合并后的内容添加到分组后的块列表

    return grouped_blocks  # 返回分组后的文本块列表

def convert_ppt_to_pdf(ppt_path):
    """Convert a PowerPoint file to PDF using LibreOffice."""
    base_name = os.path.basename(ppt_path)
    ppt_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
    new_dir_path = os.path.abspath("vectorstore/ppt_references")
    os.makedirs(new_dir_path, exist_ok=True)
    pdf_path = os.path.join(new_dir_path, f"{ppt_name_without_ext}.pdf")
    command = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', new_dir_path, ppt_path]
    subprocess.run(command, check=True)
    return pdf_path

def convert_pdf_to_images(pdf_path):
    """Convert a PDF file to a series of images using PyMuPDF."""
    doc = fitz.open(pdf_path)
    base_name = os.path.basename(pdf_path)
    pdf_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
    new_dir_path = os.path.join(os.getcwd(), "vectorstore/ppt_references")
    os.makedirs(new_dir_path, exist_ok=True)
    image_paths = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap()
        output_image_path = os.path.join(new_dir_path, f"{pdf_name_without_ext}_{page_num:04d}.png")
        pix.save(output_image_path)
        image_paths.append((output_image_path, page_num))
    doc.close()
    return image_paths

def extract_text_and_notes_from_ppt(ppt_path):
    """Extract text and notes from a PowerPoint file."""
    prs = Presentation(ppt_path)
    text_and_notes = []
    for slide in prs.slides:
        slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
        except:
            notes = ''
        text_and_notes.append((slide_text, notes))
    return text_and_notes

def save_uploaded_file(uploaded_file):
    """Save an uploaded file to a temporary directory."""
    temp_dir = os.path.join(os.getcwd(), "vectorstore", "ppt_references", "tmp")
    os.makedirs(temp_dir, exist_ok=True)
    temp_file_path = os.path.join(temp_dir, uploaded_file.name)

    with open(temp_file_path, "wb") as temp_file:
        temp_file.write(uploaded_file.read())

    return temp_file_path
